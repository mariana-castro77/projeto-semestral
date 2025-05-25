from pptx import Presentation

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]
title.text = "Planning Poker: Técnica de Estimativa em Equipes Ágeis"
subtitle.text = "Uma introdução ao uso do Planning Poker para estimativas em equipes Scrum"

# Slide 2: O que é o Planning Poker?
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_2.shapes.title
content = slide_2.shapes.placeholders[1]
title.text = "O que é o Planning Poker?"
content.text = ("- Definição:\n"
                "  Técnica de estimativa colaborativa usada por equipes ágeis para avaliar o esforço necessário para concluir uma tarefa ou história de usuário.\n\n"
                "- Objetivo:\n"
                "  Estimar de maneira coletiva e melhorar a precisão das previsões.")

# Slide 3: Como Funciona o Planning Poker
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_3.shapes.title
content = slide_3.shapes.placeholders[1]
title.text = "Como Funciona o Planning Poker"
content.text = ("- Cada membro da equipe escolhe uma carta com uma estimativa de esforço para a tarefa.\n"
                "- As cartas variam de 1 a 100 pontos, geralmente com valores de Fibonacci (1, 2, 3, 5, 8, 13, 21, 34...).\n"
                "- O facilitador apresenta uma história de usuário ou tarefa.\n"
                "- Cada membro escolhe uma carta em segredo e a revela ao mesmo tempo.\n"
                "- Caso haja discordâncias, a equipe discute as razões para as diferenças.\n"
                "- O processo é repetido até que um consenso seja alcançado.")

# Slide 4: Benefícios do Planning Poker
slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_4.shapes.title
content = slide_4.shapes.placeholders[1]
title.text = "Benefícios do Planning Poker"
content.text = ("- Estimativas colaborativas: Envolve todos os membros da equipe, garantindo que as perspectivas diversas sejam consideradas.\n\n"
                "- Evita viés de autoridade: Nenhum membro influencia as estimativas de outros, pois todos revelam suas cartas ao mesmo tempo.\n\n"
                "- Ajuste dinâmico: Permite reavaliações contínuas à medida que a equipe se torna mais experiente.\n\n"
                "- Aumento de comunicação: Promove discussões sobre requisitos e possíveis desafios.")

# Slide 5: Vantagens e Desvantagens
slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_5.shapes.title
content = slide_5.shapes.placeholders[1]
title.text = "Vantagens e Desvantagens"
content.text = ("- Vantagens:\n"
                "  - Facilidade de uso.\n"
                "  - Engajamento da equipe.\n"
                "  - Estimativas mais precisas com a colaboração.\n\n"
                "- Desvantagens:\n"
                "  - Pode ser demorado se houver muita discordância.\n"
                "  - Requer facilitação experiente.")

# Slide 6: Exemplos de Uso do Planning Poker
slide_6 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_6.shapes.title
content = slide_6.shapes.placeholders[1]
title.text = "Exemplos de Uso do Planning Poker"
content.text = ("- Equipes Scrum usando o Planning Poker em sessões de refinamento de backlog.\n"
                "- Uso em equipes de desenvolvimento de software, marketing e outros projetos que necessitam de estimativas de tarefas complexas.")

# Slide 7: Dicas para uma Sessão de Planning Poker Eficiente
slide_7 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_7.shapes.title
content = slide_7.shapes.placeholders[1]
title.text = "Dicas para uma Sessão de Planning Poker Eficiente"
content.text = ("- Prepare-se bem: Antes da sessão, garanta que todos os participantes compreendam os requisitos da história de usuário.\n\n"
                "- Evite influências externas: Não deixe que discussões externas ou a opinião de um único membro influenciem a estimativa.\n\n"
                "- Respeite as diferenças: Entenda que as diferentes estimativas podem indicar diferentes pontos de vista sobre a tarefa.")

# Slide 8: Conclusão
slide_8 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_8.shapes.title
content = slide_8.shapes.placeholders[1]
title.text = "Conclusão"
content.text = ("- O Planning Poker é uma técnica eficaz para estimativas em ambientes ágeis, incentivando a colaboração e o consenso dentro da equipe.\n"
                "- Ao aplicar o Planning Poker, equipes podem melhorar a precisão das estimativas e a comunicação entre os membros.")

# Slide 9: Perguntas e Discussão
slide_9 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_9.shapes.title
content = slide_9.shapes.placeholders[1]
title.text = "Perguntas e Discussão"
content.text = ("- Alguma dúvida?\n"
                "- Como você aplicaria o Planning Poker em sua equipe?")

# Save the presentation
pptx_path = '/mnt/data/planning_poker_presentation.pptx'
prs.save(pptx_path)

pptx_path
